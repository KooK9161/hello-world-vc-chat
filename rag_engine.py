"""
RAG Engine - Document chunking, embedding, and retrieval using ChromaDB
Supports: .txt, .md, .csv, .json, .docx, .pdf, .xlsx, .pptx
"""

import os
import re
import chromadb


def extract_text_from_file(filepath):
    """Standalone text extraction from any supported file format.
    Returns extracted text as string. Used by chat file attachment feature.
    """
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext in ('.txt', '.md', '.csv'):
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()

        elif ext == '.json':
            import json
            with open(filepath, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return json.dumps(data, ensure_ascii=False, indent=2)

        elif ext == '.docx':
            import docx
            doc = docx.Document(filepath)
            parts = []
            for p in doc.paragraphs:
                if p.text.strip():
                    parts.append(p.text)
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(' | '.join(cells))
            return '\n\n'.join(parts)

        elif ext == '.pdf':
            from PyPDF2 import PdfReader
            reader = PdfReader(filepath)
            pages = [page.extract_text() or '' for page in reader.pages]
            return '\n\n'.join(pages)

        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(filepath)
            parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = [f"[Slide {slide_num}]"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_texts.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                parts.append(' | '.join(cells))
                if len(slide_texts) > 1:
                    parts.append('\n'.join(slide_texts))
            return '\n\n'.join(parts)

        elif ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            text_parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text_parts.append(f"[Sheet: {sheet}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else '' for c in row]
                    line = '\t'.join(cells).strip()
                    if line and line != '\t' * len(cells):
                        text_parts.append(line)
            wb.close()
            return '\n'.join(text_parts)

    except Exception as e:
        print(f"[Extract] Error reading {filepath}: {e}")
        return ''

    return ''


class RAGEngine:
    """RAG engine for a single persona's knowledge base."""

    def __init__(self, persona_id, knowledge_path, db_path):
        self.persona_id = persona_id
        self.knowledge_path = knowledge_path
        self.collection_name = f"persona_{persona_id}"

        # Persistent ChromaDB client
        self.client = chromadb.PersistentClient(path=db_path)
        self.collection = self.client.get_or_create_collection(
            name=self.collection_name,
            metadata={"hnsw:space": "cosine"}
        )

        # Indexing status
        self.indexing = False
        self.index_error = None

    def _read_file(self, filepath):
        """Read file content based on extension."""
        ext = os.path.splitext(filepath)[1].lower()

        try:
            if ext in ('.txt', '.md', '.csv'):
                with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                    return f.read()

            elif ext == '.json':
                import json
                with open(filepath, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    return json.dumps(data, ensure_ascii=False, indent=2)

            elif ext == '.docx':
                try:
                    import docx
                    doc = docx.Document(filepath)
                    parts = []
                    for p in doc.paragraphs:
                        if p.text.strip():
                            parts.append(p.text)
                    # Also extract tables
                    for table in doc.tables:
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                parts.append(' | '.join(cells))
                    return '\n\n'.join(parts)
                except ImportError:
                    print("[RAG] python-docx not installed")
                    return ''

            elif ext == '.pdf':
                try:
                    from PyPDF2 import PdfReader
                    reader = PdfReader(filepath)
                    pages = [page.extract_text() or '' for page in reader.pages]
                    return '\n\n'.join(pages)
                except ImportError:
                    print("[RAG] PyPDF2 not installed")
                    return ''

            elif ext == '.pptx':
                try:
                    from pptx import Presentation
                    prs = Presentation(filepath)
                    parts = []
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_texts = [f"[Slide {slide_num}]"]
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    text = para.text.strip()
                                    if text:
                                        slide_texts.append(text)
                            if shape.has_table:
                                for row in shape.table.rows:
                                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                                    if cells:
                                        slide_texts.append(' | '.join(cells))
                        if len(slide_texts) > 1:  # More than just the header
                            parts.append('\n'.join(slide_texts))
                    return '\n\n'.join(parts)
                except ImportError:
                    print("[RAG] python-pptx not installed")
                    return ''

            elif ext == '.xlsx':
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
                    text_parts = []
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        text_parts.append(f"[Sheet: {sheet}]")
                        for row in ws.iter_rows(values_only=True):
                            cells = [str(c) if c is not None else '' for c in row]
                            line = '\t'.join(cells).strip()
                            if line and line != '\t' * len(cells):
                                text_parts.append(line)
                    wb.close()
                    return '\n'.join(text_parts)
                except ImportError:
                    print("[RAG] openpyxl not installed")
                    return ''

        except Exception as e:
            print(f"[RAG] Error reading {filepath}: {e}")
            return ''

        return ''

    def _chunk_text(self, text, chunk_size=400, overlap=80):
        """Split text into overlapping chunks by sentences."""
        if not text.strip():
            return []

        # Split by paragraphs first, then by sentences
        sentences = re.split(r'(?<=[.!?。！？\n])\s*', text)
        sentences = [s.strip() for s in sentences if s.strip()]

        chunks = []
        current_chunk = []
        current_length = 0

        for sentence in sentences:
            words = sentence.split()
            sentence_length = len(words)

            if current_length + sentence_length > chunk_size and current_chunk:
                chunk_text = ' '.join(' '.join(s.split()) for s in current_chunk)
                if chunk_text.strip():
                    chunks.append(chunk_text)

                # Keep overlap
                overlap_sentences = []
                overlap_length = 0
                for s in reversed(current_chunk):
                    s_len = len(s.split())
                    if overlap_length + s_len > overlap:
                        break
                    overlap_sentences.insert(0, s)
                    overlap_length += s_len

                current_chunk = overlap_sentences
                current_length = overlap_length

            current_chunk.append(sentence)
            current_length += sentence_length

        # Last chunk
        if current_chunk:
            chunk_text = ' '.join(' '.join(s.split()) for s in current_chunk)
            if chunk_text.strip():
                chunks.append(chunk_text)

        return chunks

    def index_documents(self):
        """Index all documents in the knowledge directory."""
        self.indexing = True
        self.index_error = None

        try:
            # Delete existing collection and recreate
            try:
                self.client.delete_collection(self.collection_name)
            except Exception:
                pass
            self.collection = self.client.create_collection(
                name=self.collection_name,
                metadata={"hnsw:space": "cosine"}
            )

            all_documents = []
            all_ids = []
            all_metadatas = []

            if not os.path.exists(self.knowledge_path):
                self.indexing = False
                return

            for filename in sorted(os.listdir(self.knowledge_path)):
                filepath = os.path.join(self.knowledge_path, filename)
                if not os.path.isfile(filepath):
                    continue

                print(f"[RAG] Reading: {filename}")
                text = self._read_file(filepath)
                if not text:
                    print(f"[RAG] No text extracted from: {filename}")
                    continue

                chunks = self._chunk_text(text)
                print(f"[RAG] {filename} → {len(chunks)} chunks")

                for i, chunk in enumerate(chunks):
                    doc_id = f"{filename}_chunk_{i}"
                    all_documents.append(chunk)
                    all_ids.append(doc_id)
                    all_metadatas.append({
                        'source': filename,
                        'chunk_index': i,
                        'total_chunks': len(chunks)
                    })

            if all_documents:
                batch_size = 100
                for start in range(0, len(all_documents), batch_size):
                    end = min(start + batch_size, len(all_documents))
                    self.collection.add(
                        documents=all_documents[start:end],
                        ids=all_ids[start:end],
                        metadatas=all_metadatas[start:end]
                    )

            print(f"[RAG] Indexed {len(all_documents)} chunks total from {self.knowledge_path}")

        except Exception as e:
            self.index_error = str(e)
            print(f"[RAG] Indexing error: {e}")
        finally:
            self.indexing = False

    def retrieve(self, query, top_k=5):
        """Retrieve relevant chunks for a query."""
        if self.collection.count() == 0:
            return []

        n_results = min(top_k, self.collection.count())

        results = self.collection.query(
            query_texts=[query],
            n_results=n_results
        )

        contexts = []
        if results and results['documents'] and results['documents'][0]:
            for doc, meta in zip(results['documents'][0], results['metadatas'][0]):
                contexts.append({
                    'text': doc,
                    'source': meta.get('source', 'unknown')
                })

        return contexts

    def get_chunk_count(self):
        """Get the number of indexed chunks."""
        return self.collection.count()

    def is_indexing(self):
        """Check if indexing is in progress."""
        return self.indexing
